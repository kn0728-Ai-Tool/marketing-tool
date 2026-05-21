# digest_creator.py
# =====================================
# 複数資料のダイジェスト化 & プレゼン資料生成モジュール
# 対応入力: PDF / Word(.docx) / テキスト
# 対応出力: PowerPoint(.pptx) / Excel(.xlsx)
# =====================================

import io
import json
import re
import datetime
from zoneinfo import ZoneInfo

JST = ZoneInfo("Asia/Tokyo")


# =====================================
# テキスト抽出
# =====================================

def extract_text_from_file(uploaded_file) -> tuple[str, str]:
    """
    アップロードされたファイルからテキストを抽出する。
    戻り値: (抽出テキスト, エラーメッセージ)
    """
    name = uploaded_file.name
    ext  = name.split(".")[-1].lower()
    data = uploaded_file.read()

    try:
        if ext == "pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(io.BytesIO(data))
                text = "\n".join(
                    page.extract_text() or "" for page in reader.pages
                )
                return text.strip(), ""
            except ImportError:
                return "", "pypdf がインストールされていません。requirements.txt に pypdf を追加してください。"

        elif ext == "docx":
            try:
                import docx
                doc = docx.Document(io.BytesIO(data))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                return text.strip(), ""
            except ImportError:
                return "", "python-docx がインストールされていません。requirements.txt に python-docx を追加してください。"

        elif ext in ("txt", "text", "md"):
            try:
                text = data.decode("utf-8")
            except UnicodeDecodeError:
                text = data.decode("shift-jis", errors="replace")
            return text.strip(), ""

        else:
            return "", f"非対応のファイル形式です: .{ext}"

    except Exception as e:
        return "", f"読み込みエラー ({name}): {e}"


# =====================================
# AIによるダイジェスト生成
# =====================================

def generate_digest(
    texts: dict,          # {ファイル名: テキスト}
    api_key: str,
    num_slides: int = 8,
    detail_level: str = "標準",
) -> dict:
    """
    複数資料を横断してダイジェスト化し、スライド構成を生成する。

    戻り値:
    {
      "title":    str,               # プレゼン全体タイトル
      "subtitle": str,               # サブタイトル
      "overview": str,               # 全体サマリー（1〜2段落）
      "slides": [
        {
          "slide_num":  int,
          "title":      str,
          "layout":     "bullets" | "two_col" | "big_number" | "summary",
          "bullets":    [str, ...],  # 箇条書きポイント
          "left_text":  str,         # 2列レイアウト左側
          "right_text": str,         # 2列レイアウト右側
          "big_number": str,         # 大数字（big_numberレイアウト用）
          "big_label":  str,         # 大数字のラベル
          "notes":      str,         # 発表者ノート
        },
        ...
      ],
      "key_points":   [str, ...],    # 全資料の重要ポイント（Excel用）
      "action_items": [str, ...],    # アクションアイテム
    }
    """
    from openai import OpenAI
    client = OpenAI(api_key=api_key)

    # 資料テキストを結合（長すぎる場合は各資料を先頭から制限）
    MAX_CHARS_PER_FILE = 3000
    combined = ""
    for fname, text in texts.items():
        truncated = text[:MAX_CHARS_PER_FILE] + ("..." if len(text) > MAX_CHARS_PER_FILE else "")
        combined += f"\n\n【資料：{fname}】\n{truncated}"

    detail_map = {
        "簡潔": "各スライドは箇条書き2〜3点、シンプルに",
        "標準": "各スライドは箇条書き3〜4点、適度な詳しさで",
        "詳細": "各スライドは箇条書き4〜5点、詳しく丁寧に",
    }
    detail_instruction = detail_map.get(detail_level, detail_map["標準"])

    prompt = f"""
あなたはプロのプレゼンテーション作成者です。
以下の複数の資料を読み込み、プレゼン資料のスライド構成をJSON形式で作成してください。

【条件】
- スライド枚数: タイトルスライド1枚 + コンテンツ{num_slides - 2}枚 + まとめスライド1枚 = 合計{num_slides}枚
- {detail_instruction}
- 日本語で出力すること
- JSON以外のテキストは一切出力しないこと

【レイアウト種別】
- "bullets"     : 見出し＋箇条書き（最も一般的）
- "two_col"     : 左右2列（比較・対比に使用）
- "big_number"  : 大きな数字＋説明（統計・指標の強調に使用）
- "summary"     : まとめスライド（最終スライドに使用）

【出力JSONフォーマット】
{{
  "title": "プレゼン全体のタイトル",
  "subtitle": "サブタイトルまたは日付・部署名など",
  "overview": "全資料を横断した要約（2〜3文）",
  "slides": [
    {{
      "slide_num": 1,
      "title": "スライドタイトル",
      "layout": "bullets",
      "bullets": ["ポイント1", "ポイント2", "ポイント3"],
      "left_text": "",
      "right_text": "",
      "big_number": "",
      "big_label": "",
      "notes": "発表者向けの補足メモ"
    }}
  ],
  "key_points": ["重要ポイント1", "重要ポイント2", "重要ポイント3", "重要ポイント4", "重要ポイント5"],
  "action_items": ["アクション1", "アクション2", "アクション3"]
}}

【入力資料】
{combined}
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "あなたはプロのプレゼンテーション作成の専門家です。"},
            {"role": "user",   "content": prompt},
        ],
        temperature=0.4,
        max_tokens=3000,
    )

    raw = response.choices[0].message.content.strip()
    json_match = re.search(r"\{.*\}", raw, re.DOTALL)
    if not json_match:
        return {"error": "AIの応答からJSONを取得できませんでした。", "raw": raw}

    try:
        return json.loads(json_match.group())
    except json.JSONDecodeError as e:
        return {"error": f"JSONパースエラー: {e}", "raw": raw}


# =====================================
# PowerPoint生成
# =====================================

def create_pptx(digest: dict) -> bytes:
    """
    ダイジェストデータからPowerPointファイルを生成する。
    戻り値: pptxファイルのバイトデータ
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # カラーパレット（Midnight Executive）
    COLOR_BG_DARK   = RGBColor(0x1E, 0x27, 0x61)   # ネイビー（タイトル・まとめ背景）
    COLOR_BG_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # 白（コンテンツ背景）
    COLOR_ACCENT    = RGBColor(0x63, 0x66, 0xF1)   # インディゴ（アクセント）
    COLOR_TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)   # ほぼ黒（本文）
    COLOR_TEXT_LIGHT= RGBColor(0xFF, 0xFF, 0xFF)   # 白（ダーク背景上）
    COLOR_TEXT_SUB  = RGBColor(0x64, 0x74, 0x8B)   # グレー（サブテキスト）
    COLOR_BULLET_BG = RGBColor(0xCA, 0xDC, 0xFC)   # アイスブルー（箇条書き背景）

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 完全空白レイアウト

    def add_rect(slide, x, y, w, h, color, transparency=0):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if transparency > 0:
            shape.fill.fore_color.theme_color = None
        shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, size, color, bold=False,
                 align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = "Calibri"
        return txBox

    def add_bullet_box(slide, bullets, x, y, w, h, font_size=14):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet_text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = bullet_text
            run.font.size = Pt(font_size)
            run.font.color.rgb = COLOR_TEXT_DARK
            run.font.name = "Calibri"
            # 手動で bullet 記号を先頭に追加
            run.text = f"▶  {bullet_text}"
        return txBox

    slides_data = digest.get("slides", [])

    for slide_data in slides_data:
        slide  = prs.slides.add_slide(blank_layout)
        layout = slide_data.get("layout", "bullets")
        title  = slide_data.get("title", "")
        num    = slide_data.get("slide_num", 1)

        # ── タイトルスライド（1枚目） ──
        if num == 1:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG_DARK

            # アクセントバー（左側）
            add_rect(slide, 0, 0, 0.5, 7.5, COLOR_ACCENT)

            # タイトル
            add_text(slide, digest.get("title", title),
                     1.0, 2.2, 11.0, 1.4,
                     size=40, color=COLOR_TEXT_LIGHT, bold=True,
                     align=PP_ALIGN.LEFT)
            # サブタイトル
            add_text(slide, digest.get("subtitle", ""),
                     1.0, 3.8, 11.0, 0.7,
                     size=18, color=COLOR_BULLET_BG,
                     align=PP_ALIGN.LEFT)
            # 日付
            now_str = datetime.datetime.now(JST).strftime("%Y年%m月%d日")
            add_text(slide, now_str,
                     1.0, 6.5, 6.0, 0.5,
                     size=12, color=COLOR_TEXT_SUB,
                     align=PP_ALIGN.LEFT)

        # ── まとめスライド ──
        elif layout == "summary":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG_DARK
            add_rect(slide, 0, 0, 0.5, 7.5, COLOR_ACCENT)

            add_text(slide, title,
                     1.0, 0.4, 11.5, 0.9,
                     size=32, color=COLOR_TEXT_LIGHT, bold=True)

            bullets = slide_data.get("bullets", [])
            for i, pt in enumerate(bullets[:5]):
                y_pos = 1.5 + i * 0.95
                add_rect(slide, 1.0, y_pos, 11.0, 0.75, COLOR_ACCENT)
                add_text(slide, f"  {pt}",
                         1.0, y_pos + 0.05, 11.0, 0.65,
                         size=14, color=COLOR_TEXT_LIGHT)

        # ── big_number レイアウト ──
        elif layout == "big_number":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG_LIGHT
            add_rect(slide, 0, 0, 13.33, 1.1, COLOR_ACCENT)

            add_text(slide, title,
                     0.5, 0.15, 12.3, 0.8,
                     size=26, color=COLOR_TEXT_LIGHT, bold=True)

            add_text(slide, slide_data.get("big_number", ""),
                     3.5, 1.5, 6.0, 2.5,
                     size=80, color=COLOR_ACCENT, bold=True,
                     align=PP_ALIGN.CENTER)

            add_text(slide, slide_data.get("big_label", ""),
                     2.0, 4.1, 9.0, 0.8,
                     size=18, color=COLOR_TEXT_SUB,
                     align=PP_ALIGN.CENTER)

            bullets = slide_data.get("bullets", [])
            if bullets:
                add_bullet_box(slide, bullets, 0.6, 5.1, 12.0, 2.0, font_size=13)

        # ── two_col レイアウト ──
        elif layout == "two_col":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG_LIGHT
            add_rect(slide, 0, 0, 13.33, 1.1, COLOR_ACCENT)

            add_text(slide, title,
                     0.5, 0.15, 12.3, 0.8,
                     size=26, color=COLOR_TEXT_LIGHT, bold=True)

            # 左列
            add_rect(slide, 0.5, 1.3, 5.8, 5.5, RGBColor(0xF1, 0xF5, 0xF9))
            left_text = slide_data.get("left_text", "")
            add_text(slide, left_text,
                     0.7, 1.5, 5.4, 5.1,
                     size=13, color=COLOR_TEXT_DARK, word_wrap=True)

            # 右列
            add_rect(slide, 6.9, 1.3, 5.9, 5.5, RGBColor(0xEE, 0xF2, 0xFF))
            right_text = slide_data.get("right_text", "")
            add_text(slide, right_text,
                     7.1, 1.5, 5.5, 5.1,
                     size=13, color=COLOR_TEXT_DARK, word_wrap=True)

        # ── bullets レイアウト（デフォルト） ──
        else:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = COLOR_BG_LIGHT
            add_rect(slide, 0, 0, 13.33, 1.1, COLOR_ACCENT)

            add_text(slide, title,
                     0.5, 0.15, 12.3, 0.8,
                     size=26, color=COLOR_TEXT_LIGHT, bold=True)

            bullets = slide_data.get("bullets", [])
            for i, pt in enumerate(bullets[:5]):
                y_pos = 1.4 + i * 1.0
                # 背景カード
                card_color = COLOR_BULLET_BG if i % 2 == 0 else RGBColor(0xEE, 0xF2, 0xFF)
                add_rect(slide, 0.6, y_pos, 12.1, 0.82, card_color)
                # 番号バッジ
                add_rect(slide, 0.6, y_pos, 0.45, 0.82, COLOR_ACCENT)
                add_text(slide, str(i + 1),
                         0.6, y_pos + 0.18, 0.45, 0.5,
                         size=13, color=COLOR_TEXT_LIGHT, bold=True,
                         align=PP_ALIGN.CENTER)
                # テキスト
                add_text(slide, pt,
                         1.2, y_pos + 0.1, 11.3, 0.65,
                         size=14, color=COLOR_TEXT_DARK)

        # 発表者ノートを追加
        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    # スライド番号フッター（全スライド）
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        add_text(slide, f"{i + 1} / {len(prs.slides)}",
                 12.0, 7.0, 1.2, 0.4,
                 size=10, color=COLOR_TEXT_SUB,
                 align=PP_ALIGN.RIGHT)

    # バイト列として返す
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# =====================================
# Excel生成
# =====================================

def create_xlsx(digest: dict, texts: dict) -> bytes:
    """
    ダイジェストデータからExcelサマリーファイルを生成する。
    戻り値: xlsxファイルのバイトデータ
    """
    from openpyxl import Workbook
    from openpyxl.styles import (
        Font, PatternFill, Alignment, Border, Side, GradientFill
    )
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # ── シート①：全体サマリー ──
    ws1 = wb.active
    ws1.title = "全体サマリー"

    HDR_FILL  = PatternFill("solid", fgColor="6366F1")
    SUB_FILL  = PatternFill("solid", fgColor="EEF2FF")
    CARD_FILL = PatternFill("solid", fgColor="F8FAFC")
    HDR_FONT  = Font(name="Calibri", bold=True, color="FFFFFF", size=12)
    BODY_FONT = Font(name="Calibri", size=11)
    BOLD_FONT = Font(name="Calibri", bold=True, size=11)
    thin = Side(style="thin", color="E2E8F0")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def write_header(ws, row, col, text):
        cell = ws.cell(row=row, column=col, value=text)
        cell.font = HDR_FONT
        cell.fill = HDR_FILL
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
        return cell

    def write_cell(ws, row, col, text, fill=None, bold=False):
        cell = ws.cell(row=row, column=col, value=text)
        cell.font = BOLD_FONT if bold else BODY_FONT
        cell.fill = fill or CARD_FILL
        cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        cell.border = border
        return cell

    # タイトル行
    ws1.merge_cells("A1:D1")
    title_cell = ws1["A1"]
    title_cell.value = f"📊 {digest.get('title', 'ダイジェスト')}"
    title_cell.font = Font(name="Calibri", bold=True, size=16, color="1E2761")
    title_cell.alignment = Alignment(horizontal="left", vertical="center")
    ws1.row_dimensions[1].height = 36

    ws1.merge_cells("A2:D2")
    ws1["A2"].value = f"作成日時: {datetime.datetime.now(JST).strftime('%Y-%m-%d %H:%M')} (JST)　|　対象資料数: {len(texts)}件"
    ws1["A2"].font = Font(name="Calibri", size=10, color="64748B")
    ws1["A2"].alignment = Alignment(horizontal="left", vertical="center")
    ws1.row_dimensions[2].height = 20

    # 概要
    ws1.merge_cells("A3:D3")
    ws1["A3"].value = "【全体概要】"
    ws1["A3"].font = BOLD_FONT
    ws1["A3"].fill = SUB_FILL
    ws1.row_dimensions[3].height = 20

    ws1.merge_cells("A4:D4")
    ws1["A4"].value = digest.get("overview", "")
    ws1["A4"].font = BODY_FONT
    ws1["A4"].alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
    ws1.row_dimensions[4].height = 60

    # 重要ポイント
    ws1.merge_cells("A6:D6")
    ws1["A6"].value = "【重要ポイント】"
    ws1["A6"].font = BOLD_FONT
    ws1["A6"].fill = SUB_FILL
    ws1.row_dimensions[6].height = 22

    write_header(ws1, 7, 1, "No.")
    write_header(ws1, 7, 2, "重要ポイント")
    ws1.merge_cells("B7:D7")

    for i, kp in enumerate(digest.get("key_points", []), start=1):
        row = 7 + i
        write_cell(ws1, row, 1, i, bold=True)
        c = ws1.cell(row=row, column=2, value=kp)
        c.font = BODY_FONT
        c.fill = CARD_FILL
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        c.border = border
        ws1.merge_cells(f"B{row}:D{row}")
        ws1.row_dimensions[row].height = 30

    # アクションアイテム
    action_start = 7 + len(digest.get("key_points", [])) + 2
    ws1.merge_cells(f"A{action_start}:D{action_start}")
    ws1[f"A{action_start}"].value = "【アクションアイテム】"
    ws1[f"A{action_start}"].font = BOLD_FONT
    ws1[f"A{action_start}"].fill = SUB_FILL

    write_header(ws1, action_start + 1, 1, "No.")
    write_header(ws1, action_start + 1, 2, "アクション内容")
    write_header(ws1, action_start + 1, 3, "担当者")
    write_header(ws1, action_start + 1, 4, "期限")
    ws1.merge_cells(f"B{action_start + 1}:B{action_start + 1}")

    for i, action in enumerate(digest.get("action_items", []), start=1):
        row = action_start + 1 + i
        write_cell(ws1, row, 1, i, bold=True)
        c = ws1.cell(row=row, column=2, value=action)
        c.font = BODY_FONT
        c.fill = CARD_FILL
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        c.border = border
        write_cell(ws1, row, 3, "")
        write_cell(ws1, row, 4, "")
        ws1.row_dimensions[row].height = 28

    # 列幅
    ws1.column_dimensions["A"].width = 8
    ws1.column_dimensions["B"].width = 50
    ws1.column_dimensions["C"].width = 18
    ws1.column_dimensions["D"].width = 18

    # ── シート②：スライド構成一覧 ──
    ws2 = wb.create_sheet("スライド構成")

    write_header(ws2, 1, 1, "No.")
    write_header(ws2, 1, 2, "スライドタイトル")
    write_header(ws2, 1, 3, "レイアウト")
    write_header(ws2, 1, 4, "主なポイント")
    write_header(ws2, 1, 5, "発表者メモ")
    ws2.row_dimensions[1].height = 24

    for sd in digest.get("slides", []):
        row = sd["slide_num"] + 1
        write_cell(ws2, row, 1, sd["slide_num"], bold=True)
        write_cell(ws2, row, 2, sd.get("title", ""))
        write_cell(ws2, row, 3, sd.get("layout", ""))
        bullets = sd.get("bullets", [])
        write_cell(ws2, row, 4, "\n".join(f"・{b}" for b in bullets))
        write_cell(ws2, row, 5, sd.get("notes", ""))
        ws2.row_dimensions[row].height = max(30, len(bullets) * 18)

    ws2.column_dimensions["A"].width = 6
    ws2.column_dimensions["B"].width = 30
    ws2.column_dimensions["C"].width = 14
    ws2.column_dimensions["D"].width = 50
    ws2.column_dimensions["E"].width = 40

    # ── シート③：読み込み資料一覧 ──
    ws3 = wb.create_sheet("読み込み資料")
    write_header(ws3, 1, 1, "No.")
    write_header(ws3, 1, 2, "ファイル名")
    write_header(ws3, 1, 3, "文字数")
    write_header(ws3, 1, 4, "先頭200文字プレビュー")

    for i, (fname, text) in enumerate(texts.items(), start=1):
        row = i + 1
        write_cell(ws3, row, 1, i, bold=True)
        write_cell(ws3, row, 2, fname)
        write_cell(ws3, row, 3, len(text))
        preview = text[:200].replace("\n", " ") + ("..." if len(text) > 200 else "")
        c = ws3.cell(row=row, column=4, value=preview)
        c.font = BODY_FONT
        c.fill = CARD_FILL
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        c.border = border
        ws3.row_dimensions[row].height = 40

    ws3.column_dimensions["A"].width = 6
    ws3.column_dimensions["B"].width = 30
    ws3.column_dimensions["C"].width = 10
    ws3.column_dimensions["D"].width = 60

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()
